"""
pptx_utils.py
=============
Generates publication-quality, 16:9 widescreen Microsoft PowerPoint (.pptx)
examination presentations with embedded Matplotlib diagrams, standard mathematical
formula rendering, subject-aware layouts (Mathematics, History, Sciences, etc.),
blueprint overview tables, and marking scheme notes.
"""

import os
import re
from typing import List, Dict, Any, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# ── Color Palette (Refined Academic & Executive Theme) ──────────────────────
C_NAVY_DARK  = RGBColor(15, 23, 42)      # #0f172a - Primary dark header
C_NAVY_LIGHT = RGBColor(30, 41, 59)      # #1e293b - Secondary dark
C_BLUE_ROYAL = RGBColor(30, 64, 175)     # #1e40af - Accent blue
C_BLUE_SKY   = RGBColor(2, 132, 199)     # #0284c7 - Marks badge / highlight
C_BLUE_LIGHT = RGBColor(239, 246, 255)   # #eff6ff - Shaded card background
C_SLATE_GRAY = RGBColor(71, 85, 105)     # #475569 - Subtitles & metadata
C_BODY_DARK  = RGBColor(15, 23, 42)      # #0f172a - Main question text
C_BORDER_GRAY= RGBColor(203, 213, 225)   # #cbd5e1 - Card borders
C_WHITE      = RGBColor(255, 255, 255)   # #ffffff
C_AMBER      = RGBColor(217, 119, 6)     # #d97706 - History / Humanities badge
C_EMERALD    = RGBColor(5, 150, 105)     # #059669 - Science / Bio badge
C_PURPLE     = RGBColor(124, 58, 237)    # #7c3aed - CS badge
C_CARD_BG    = RGBColor(248, 250, 252)   # #f8fafc - Shaded card fill


# ── Subject Icons & Colors ──────────────────────────────────────────────────
SUBJECT_THEME_MAP = {
    "mathematics": {"badge": "📐 MATHEMATICS", "color": C_BLUE_ROYAL},
    "history":     {"badge": "📜 HISTORY & HUMANITIES", "color": C_AMBER},
    "physics":     {"badge": "🔬 PHYSICS", "color": C_BLUE_SKY},
    "chemistry":   {"badge": "🧪 CHEMISTRY", "color": C_EMERALD},
    "biology":     {"badge": "🧬 BIOLOGY", "color": C_EMERALD},
    "computer science": {"badge": "💻 COMPUTER SCIENCE", "color": C_PURPLE},
    "economics":   {"badge": "📊 ECONOMICS", "color": C_BLUE_SKY},
    "geography":   {"badge": "🌍 GEOGRAPHY", "color": C_EMERALD},
    "literature":  {"badge": "📚 LITERATURE", "color": C_AMBER},
    "statistics":  {"badge": "📈 STATISTICS", "color": C_BLUE_ROYAL},
}


# ── Mathematical Unicode & Clean Formatting ────────────────────────────────
def format_math_and_text(raw_text: str) -> str:
    """
    Converts LaTeX mathematical markup and raw symbols into standard,
    clean, publication-quality Unicode text suitable for presentation slides.
    """
    if not raw_text:
        return ""

    text = raw_text

    # Strip raw figure/image commands
    text = re.sub(r"\\begin\{figure\}[\s\S]*?\\end\{figure\}", "", text)
    text = re.sub(r"\\includegraphics(\[.*?\])?\{.*?\}", "", text)
    text = re.sub(r"\\caption\{.*?\}", "", text)
    text = re.sub(r"\\centering", "", text)
    text = re.sub(r"\\label\{.*?\}", "", text)

    # Fractions
    text = re.sub(r"\\frac\{([^}]+)\}\{([^}]+)\}", r"(\1 / \2)", text)
    text = re.sub(r"\\dfrac\{([^}]+)\}\{([^}]+)\}", r"(\1 / \2)", text)
    text = re.sub(r"\\tfrac\{([^}]+)\}\{([^}]+)\}", r"(\1 / \2)", text)

    # Square roots & nth roots
    text = re.sub(r"\\sqrt\[(\d+)\]\{([^}]+)\}", r"^\1√(\2)", text)
    text = re.sub(r"\\sqrt\{([^}]+)\}", r"√(\1)", text)
    text = re.sub(r"\\sqrt", "√", text)

    # Derivatives & Partials
    text = re.sub(r"\\frac\{\\partial\^?(\d+)?\s*([^}]+)\}\{\\partial\s*([^}]+)\^?(\d+)?\}", r"∂\1\2/∂\3\4", text)
    text = re.sub(r"\\partial", "∂", text)
    text = re.sub(r"\\nabla", "∇", text)

    # Integrals & Summations
    text = re.sub(r"\\int_\{([^}]+)\}\^\{([^}]+)\}", r"∫[\1 to \2] ", text)
    text = re.sub(r"\\int_([^\^ ]+)\^([^\s]+)", r"∫[\1 to \2] ", text)
    text = re.sub(r"\\iint", "∬", text)
    text = re.sub(r"\\iiint", "∭", text)
    text = re.sub(r"\\oint", "∮", text)
    text = re.sub(r"\\int", "∫", text)
    text = re.sub(r"\\sum_\{([^}]+)\}\^\{([^}]+)\}", r"∑[\1 to \2] ", text)
    text = re.sub(r"\\sum", "∑", text)
    text = re.sub(r"\\prod_\{([^}]+)\}\^\{([^}]+)\}", r"∏[\1 to \2] ", text)
    text = re.sub(r"\\prod", "∏", text)
    text = re.sub(r"\\lim_\{([^}]+)\}", r"lim(\1) ", text)

    # Greek letters
    greek_replacements = [
        (r"\\alpha", "α"), (r"\\beta", "β"), (r"\\gamma", "γ"), (r"\\delta", "δ"),
        (r"\\epsilon", "ε"), (r"\\varepsilon", "ε"), (r"\\zeta", "ζ"), (r"\\eta", "η"),
        (r"\\theta", "θ"), (r"\\vartheta", "θ"), (r"\\iota", "ι"), (r"\\kappa", "κ"),
        (r"\\lambda", "λ"), (r"\\mu", "μ"), (r"\\nu", "ν"), (r"\\xi", "ξ"),
        (r"\\pi", "π"), (r"\\varpi", "ϖ"), (r"\\rho", "ρ"), (r"\\varrho", "ϱ"),
        (r"\\sigma", "σ"), (r"\\varsigma", "ς"), (r"\\tau", "τ"), (r"\\upsilon", "υ"),
        (r"\\phi", "φ"), (r"\\varphi", "ϕ"), (r"\\chi", "χ"), (r"\\psi", "ψ"), (r"\\omega", "ω"),
        (r"\\Gamma", "Γ"), (r"\\Delta", "Δ"), (r"\\Theta", "Θ"), (r"\\Lambda", "Λ"),
        (r"\\Xi", "Ξ"), (r"\\Pi", "Π"), (r"\\Sigma", "Σ"), (r"\\Upsilon", "Υ"),
        (r"\\Phi", "Φ"), (r"\\Psi", "Ψ"), (r"\\Omega", "Ω")
    ]
    for pattern, rep in greek_replacements:
        text = re.sub(pattern, rep, text)

    # Math Symbols & Operators
    symbol_replacements = [
        (r"\\infty", "∞"), (r"\\pm", "±"), (r"\\mp", "∓"),
        (r"\\times", "×"), (r"\\div", "÷"), (r"\\cdot", "·"), (r"\\ast", "*"),
        (r"\\leq", "≤"), (r"\\le", "≤"), (r"\\geq", "≥"), (r"\\ge", "≥"),
        (r"\\neq", "≠"), (r"\\ne", "≠"), (r"\\approx", "≈"), (r"\\equiv", "≡"),
        (r"\\sim", "~"), (r"\\propto", "∝"), (r"\\to", "→"), (r"\\rightarrow", "→"),
        (r"\\leftarrow", "←"), (r"\\Rightarrow", "⇒"), (r"\\Leftarrow", "⇐"),
        (r"\\Leftrightarrow", "⇔"), (r"\\iff", "⇔"), (r"\\implies", "⇒"),
        (r"\\forall", "∀"), (r"\\exists", "∃"), (r"\\in", "∈"), (r"\\notin", "∉"),
        (r"\\subset", "⊂"), (r"\\subseteq", "⊆"), (r"\\cup", "∪"), (r"\\cap", "∩"),
        (r"\\perp", "⊥"), (r"\\parallel", "∥"), (r"\\angle", "∠"), (r"\\degree", "°"),
        (r"\\circ", "°"), (r"\\hbar", "ħ"), (r"\\ell", "ℓ"),
    ]
    for pattern, rep in symbol_replacements:
        text = re.sub(pattern, rep, text)

    # Function names
    text = re.sub(r"\\(sin|cos|tan|cot|sec|csc|arcsin|arccos|arctan|sinh|cosh|tanh|ln|log|exp|det|dim|ker|deg|max|min)\b", r"\1", text)

    # Formatting tags
    text = re.sub(r"\\mathbf\{([^}]+)\}", r"\1", text)
    text = re.sub(r"\\mathbf\s+([A-Za-z0-9])", r"\1", text)
    text = re.sub(r"\\mathbb\{R\}\^?(\d+)?", r"ℝ\1", text)
    text = re.sub(r"\\mathbb\{C\}", "ℂ", text)
    text = re.sub(r"\\mathbb\{Z\}", "ℤ", text)
    text = re.sub(r"\\mathbb\{N\}", "ℕ", text)
    text = re.sub(r"\\mathbb\{([^}]+)\}", r"\1", text)
    text = re.sub(r"\\textbf\{([^}]+)\}", r"\1", text)
    text = re.sub(r"\\textit\{([^}]+)\}", r"\1", text)
    text = re.sub(r"\\textrm\{([^}]+)\}", r"\1", text)
    text = re.sub(r"\\text\{([^}]+)\}", r"\1", text)
    text = re.sub(r"\\mathrm\{([^}]+)\}", r"\1", text)
    text = re.sub(r"\\left\(", "(", text)
    text = re.sub(r"\\right\)", ")", text)
    text = re.sub(r"\\left\[", "[", text)
    text = re.sub(r"\\right\]", "]", text)
    text = re.sub(r"\\left\\\{", "{", text)
    text = re.sub(r"\\right\\\}", "}", text)
    text = re.sub(r"\\left\|", "|", text)
    text = re.sub(r"\\right\|", "|", text)

    # Superscripts & Subscripts common in Math & Physics
    supers = {"0": "⁰", "1": "¹", "2": "²", "3": "³", "4": "⁴", "5": "⁵", "6": "⁶", "7": "⁷", "8": "⁸", "9": "⁹", "+": "⁺", "-": "⁻", "n": "ⁿ", "x": "ˣ", "t": "ᵗ"}
    for k, v in supers.items():
        text = re.sub(r"\^" + re.escape(k) + r"\b", v, text)
        text = re.sub(r"\^\{" + re.escape(k) + r"\}", v, text)

    # Clean redundant dollar signs and backslashes
    text = text.replace("$$", "").replace("$", "")
    text = re.sub(r"\\([,;! ])", r"\1", text)
    text = text.replace(r"\,", " ").replace(r"\;", " ").replace(r"\quad", "  ").replace(r"\qquad", "    ")
    text = re.sub(r"\\[a-zA-Z]+", "", text)  # remove any dangling unrecognized latex commands

    # Normalize whitespace
    text = re.sub(r"[ \t]+", " ", text).strip()
    return text


def _detect_subject_for_exam(title: str, questions: List[Dict[str, Any]]) -> str:
    """Infers the primary academic subject for the exam."""
    combined = title + " " + " ".join(q.get("text", "") for q in questions)
    from chunker import detect_subject
    return detect_subject(combined, default_subject="Mathematics")


def _add_header_banner(slide, title_text: str, badge_text: str, badge_color: RGBColor):
    """Adds a modern top banner with title and badge pill."""
    # Top banner bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_NAVY_DARK
    top_bar.line.color.rgb = C_NAVY_DARK

    # Bottom accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_BLUE_SKY
    line.line.color.rgb = C_BLUE_SKY

    # Header Title text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text[:75]
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    # Badge Pill (Top Right)
    badge_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.3), Inches(2.9), Inches(0.55))
    badge_box.fill.solid()
    badge_box.fill.fore_color.rgb = badge_color
    badge_box.line.color.rgb = badge_color
    btf = badge_box.text_frame
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.text = badge_text
    bp.font.name = "Arial"
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = C_WHITE


def _add_footer(slide, current_page: int, total_pages: int, exam_title: str):
    """Adds a subtle bottom footer bar."""
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = C_NAVY_DARK
    footer_bar.line.color.rgb = C_NAVY_DARK

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(7.08), Inches(7.5), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"ExamGen Autonomous Examination System  •  {exam_title[:45]}"
    p.font.name = "Arial"
    p.font.size = Pt(9.5)
    p.font.color.rgb = C_BORDER_GRAY

    txBox2 = slide.shapes.add_textbox(Inches(9.5), Inches(7.08), Inches(3.2), Inches(0.35))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    p2.text = f"Slide {current_page} of {total_pages}"
    p2.font.name = "Arial"
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = C_BORDER_GRAY


def create_pptx(
    exam_title: str,
    questions: List[Dict[str, Any]],
    max_marks: int = 100,
    output_path: str = "exam.pptx"
) -> str:
    """
    Builds an executive-grade, 16:9 presentation slide deck from the exam data.
    """
    prs = Presentation()
    # Set 16:9 widescreen dimensions (13.333 in x 7.5 in)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank layout

    detected_subj = _detect_subject_for_exam(exam_title, questions)
    theme = SUBJECT_THEME_MAP.get(detected_subj.lower(), {"badge": f"📌 {detected_subj.upper()}", "color": C_BLUE_ROYAL})

    total_slides = 2 + len(questions) + 1  # Title + Blueprint + Questions + Marking Scheme

    # ═════════════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE & EXAM BLUEPRINT COVER
    # ═════════════════════════════════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_layout)

    # Hero top background block
    hero = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(3.2))
    hero.fill.solid()
    hero.fill.fore_color.rgb = C_NAVY_DARK
    hero.line.color.rgb = C_NAVY_DARK

    # Accent decorative strip
    strip = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08))
    strip.fill.solid()
    strip.fill.fore_color.rgb = C_BLUE_SKY
    strip.line.color.rgb = C_BLUE_SKY

    # Subject Tag Pill
    s_badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(3.2), Inches(0.48))
    s_badge.fill.solid()
    s_badge.fill.fore_color.rgb = theme["color"]
    s_badge.line.color.rgb = theme["color"]
    s_tf = s_badge.text_frame
    sp = s_tf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.text = theme["badge"]
    sp.font.name = "Arial"
    sp.font.size = Pt(11)
    sp.font.bold = True
    sp.font.color.rgb = C_WHITE

    # Exam Title
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.4))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = exam_title.upper()
    tp.font.name = "Arial"
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = C_WHITE

    sub_p = ttf.add_paragraph()
    sub_p.text = "Official University / Board Level Comprehensive Examination Paper"
    sub_p.font.name = "Arial"
    sub_p.font.size = Pt(13)
    sub_p.font.color.rgb = C_BORDER_GRAY

    # Metadata Stat Cards (3 Column Grid)
    card_w = Inches(3.64)
    card_h = Inches(1.2)
    card_y = Inches(3.55)

    stats = [
        ("Total Max Marks", f"{max_marks} Marks", "🏆"),
        ("Question Count", f"{len(questions)} Problems", "📝"),
        ("Duration / Time", "3 Hours (180 Mins)", "⏱️"),
    ]

    for i, (label, val, icon) in enumerate(stats):
        cx = Inches(0.8) + i * (card_w + Inches(0.38))
        card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER_GRAY
        card.line.width = Pt(1.2)

        ctf = card.text_frame
        ctf.word_wrap = True
        cp1 = ctf.paragraphs[0]
        cp1.text = f"{icon}  {label}"
        cp1.font.name = "Arial"
        cp1.font.size = Pt(11)
        cp1.font.color.rgb = C_SLATE_GRAY

        cp2 = ctf.add_paragraph()
        cp2.text = val
        cp2.font.name = "Arial"
        cp2.font.size = Pt(17)
        cp2.font.bold = True
        cp2.font.color.rgb = C_NAVY_DARK

    # General Instructions Card
    inst_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8))
    inst_box.fill.solid()
    inst_box.fill.fore_color.rgb = C_BLUE_LIGHT
    inst_box.line.color.rgb = C_BLUE_ROYAL
    inst_box.line.width = Pt(1)

    itf = inst_box.text_frame
    itf.word_wrap = True
    ip0 = itf.paragraphs[0]
    ip0.text = "📋 CANDIDATE GENERAL INSTRUCTIONS:"
    ip0.font.name = "Arial"
    ip0.font.size = Pt(12)
    ip0.font.bold = True
    ip0.font.color.rgb = C_BLUE_ROYAL

    instructions = [
        "1. All questions are compulsory unless specified otherwise. Marks are indicated against each question.",
        "2. For mathematical & scientific problems, show complete step-by-step derivations, calculations, and state appropriate units.",
        "3. For humanities & history questions, provide structured source analysis, contextual evidence, and critical reasoning.",
        "4. Neat, labeled diagrams and plots must be provided wherever asked."
    ]
    for inst in instructions:
        p = itf.add_paragraph()
        p.text = inst
        p.font.name = "Arial"
        p.font.size = Pt(10.5)
        p.font.color.rgb = C_NAVY_LIGHT

    _add_footer(slide1, 1, total_slides, exam_title)

    # ═════════════════════════════════════════════════════════════════════════
    # SLIDE 2: EXAM BLUEPRINT & QUESTIONS TABLE
    # ═════════════════════════════════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_layout)
    _add_header_banner(slide2, f"Exam Blueprint & Questions Overview", theme["badge"], theme["color"])

    # Table of Questions
    rows = len(questions) + 1
    cols = 4
    table_shape = slide2.shapes.add_table(rows, cols, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.48 * rows))
    table = table_shape.table
    table.columns[0].width = Inches(1.5)   # Question #
    table.columns[1].width = Inches(6.8)   # Question Summary / Focus
    table.columns[2].width = Inches(1.8)   # Marks
    table.columns[3].width = Inches(1.6)   # Diagram / Type

    # Header Row
    headers = ["Question #", "Topic / Problem Focus", "Allocated Marks", "Visual Diagram"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY_DARK
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER if j in (0, 2, 3) else PP_ALIGN.LEFT

    # Data Rows
    for i, q in enumerate(questions, start=1):
        clean_q = format_math_and_text(q.get("text", ""))
        snippet = (clean_q[:90] + "...") if len(clean_q) > 90 else clean_q
        has_diagram = "Yes (Matplotlib)" if q.get("image_path") else "Standard"

        row_cells = [
            f"Question {i}",
            snippet,
            f"{q.get('marks', 0)} Marks",
            has_diagram
        ]
        for j, val in enumerate(row_cells):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD_BG if i % 2 == 0 else C_WHITE
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(10)
            p.font.color.rgb = C_BODY_DARK
            if j in (0, 2, 3):
                p.alignment = PP_ALIGN.CENTER
                if j == 2:
                    p.font.bold = True
                    p.font.color.rgb = C_BLUE_ROYAL

    _add_footer(slide2, 2, total_slides, exam_title)

    # ═════════════════════════════════════════════════════════════════════════
    # SLIDES 3+: INDIVIDUAL QUESTION SLIDES
    # ═════════════════════════════════════════════════════════════════════════
    for q_idx, q in enumerate(questions, start=1):
        slide = prs.slides.add_slide(blank_layout)
        marks = q.get("marks", round(max_marks / max(1, len(questions))))
        img_path = q.get("image_path")
        clean_text = format_math_and_text(q.get("text", ""))

        # Top Header Bar with Question Number
        _add_header_banner(slide, f"Question {q_idx}  [{marks} Marks]", theme["badge"], theme["color"])

        # Check if local image exists
        abs_img_path = None
        if img_path:
            # Handle relative and absolute paths
            if os.path.isabs(img_path) and os.path.exists(img_path):
                abs_img_path = img_path
            else:
                candidate = os.path.join(os.path.dirname(os.path.abspath(__file__)), img_path.lstrip("/\\"))
                if os.path.exists(candidate):
                    abs_img_path = candidate

        # Layout Variant A: With High-Resolution Embedded Plot
        if abs_img_path and os.path.exists(abs_img_path):
            # Left Column: Question Text Card
            left_w = Inches(6.8)
            left_h = Inches(5.2)
            q_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.45), left_w, left_h)
            q_card.fill.solid()
            q_card.fill.fore_color.rgb = C_CARD_BG
            q_card.line.color.rgb = C_BORDER_GRAY
            q_card.line.width = Pt(1.2)

            qtf = q_card.text_frame
            qtf.word_wrap = True
            qp0 = qtf.paragraphs[0]
            qp0.text = f"PROBLEM STATEMENT:"
            qp0.font.name = "Arial"
            qp0.font.size = Pt(11)
            qp0.font.bold = True
            qp0.font.color.rgb = C_BLUE_ROYAL

            # Split question text into paragraphs if multi-part
            paragraphs = clean_text.split("\n")
            for p_text in paragraphs:
                p_text = p_text.strip()
                if not p_text:
                    continue
                qp = qtf.add_paragraph()
                qp.text = p_text
                qp.font.name = "Arial"
                qp.font.size = Pt(14 if len(clean_text) < 300 else 12)
                qp.font.color.rgb = C_BODY_DARK

            # Right Column: Diagram Frame Card
            right_x = Inches(7.9)
            right_w = Inches(4.6)
            right_h = Inches(5.2)

            img_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(1.45), right_w, right_h)
            img_card.fill.solid()
            img_card.fill.fore_color.rgb = C_WHITE
            img_card.line.color.rgb = C_BORDER_GRAY
            img_card.line.width = Pt(1.2)

            # Insert Image inside frame
            try:
                slide.shapes.add_picture(
                    abs_img_path,
                    right_x + Inches(0.2),
                    Inches(1.65),
                    width=right_w - Inches(0.4)
                )
            except Exception as e:
                print(f"[PPTX Image Warning] Could not embed picture {abs_img_path}: {e}")

            # Caption below image
            cap_box = slide.shapes.add_textbox(right_x, Inches(6.15), right_w, Inches(0.4))
            cap_tf = cap_box.text_frame
            cap_p = cap_tf.paragraphs[0]
            cap_p.alignment = PP_ALIGN.CENTER
            cap_p.text = f"Figure {q_idx}: Scientific Plot / Diagram"
            cap_p.font.name = "Arial"
            cap_p.font.size = Pt(9.5)
            cap_p.font.italic = True
            cap_p.font.color.rgb = C_SLATE_GRAY

        # Layout Variant B: Full-Width Card (Text, Math, History Source Quotes)
        else:
            card_w = Inches(11.7)
            card_h = Inches(4.1)
            q_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.45), card_w, card_h)
            q_card.fill.solid()
            q_card.fill.fore_color.rgb = C_CARD_BG
            q_card.line.color.rgb = C_BORDER_GRAY
            q_card.line.width = Pt(1.2)

            qtf = q_card.text_frame
            qtf.word_wrap = True
            qp0 = qtf.paragraphs[0]
            qp0.text = "PROBLEM STATEMENT & TASKS:"
            qp0.font.name = "Arial"
            qp0.font.size = Pt(12)
            qp0.font.bold = True
            qp0.font.color.rgb = C_BLUE_ROYAL

            # Detect if this is a History primary source or quote question
            if '"' in clean_text or '“' in clean_text or detected_subj == "History":
                paragraphs = clean_text.split("\n")
                for p_text in paragraphs:
                    p_text = p_text.strip()
                    if not p_text:
                        continue
                    qp = qtf.add_paragraph()
                    qp.text = p_text
                    qp.font.name = "Arial"
                    qp.font.size = Pt(15 if len(clean_text) < 400 else 13)
                    qp.font.color.rgb = C_BODY_DARK
            else:
                # Standard Math / Science text with equations
                paragraphs = clean_text.split("\n")
                for p_text in paragraphs:
                    p_text = p_text.strip()
                    if not p_text:
                        continue
                    qp = qtf.add_paragraph()
                    qp.text = p_text
                    qp.font.name = "Arial"
                    qp.font.size = Pt(16 if len(clean_text) < 300 else 13.5)
                    qp.font.color.rgb = C_BODY_DARK

            # Bottom Answer Area / Evaluation Guideline Strip
            guide_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.75), card_w, Inches(0.95))
            guide_card.fill.solid()
            guide_card.fill.fore_color.rgb = C_BLUE_LIGHT
            guide_card.line.color.rgb = C_BLUE_ROYAL
            guide_card.line.width = Pt(1)

            gtf = guide_card.text_frame
            gtf.word_wrap = True
            gp0 = gtf.paragraphs[0]
            gp0.text = "💡 EVALUATION & ANSWER DIRECTIVES:"
            gp0.font.name = "Arial"
            gp0.font.size = Pt(10)
            gp0.font.bold = True
            gp0.font.color.rgb = C_BLUE_ROYAL

            gp1 = gtf.add_paragraph()
            gp1.text = f"Total Marks: {marks}  •  Candidates must demonstrate rigorous mathematical proof / textual substantiation and clearly outline final conclusions."
            gp1.font.name = "Arial"
            gp1.font.size = Pt(10)
            gp1.font.color.rgb = C_NAVY_LIGHT

        # Presenter Notes with Question Details & Marking Rubric
        slide.notes_slide.notes_text_frame.text = (
            f"Question {q_idx} Notes:\n"
            f"- Allocated Marks: {marks}\n"
            f"- Subject: {detected_subj}\n"
            f"- Text: {clean_text}\n"
            f"- Marking Guidance: Award full marks for complete and rigorous answers. Deduct marks for missing steps, wrong algebraic units, or lack of critical citations."
        )

        _add_footer(slide, 2 + q_idx, total_slides, exam_title)

    # ═════════════════════════════════════════════════════════════════════════
    # FINAL SLIDE: MARKING SCHEME & RUBRIC GUIDELINES
    # ═════════════════════════════════════════════════════════════════════════
    slide_final = prs.slides.add_slide(blank_layout)
    _add_header_banner(slide_final, "Examiner Marking Scheme & Rubric Guidelines", "GRADING RUBRIC", C_EMERALD)

    rubric_card = slide_final.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.1))
    rubric_card.fill.solid()
    rubric_card.fill.fore_color.rgb = C_CARD_BG
    rubric_card.line.color.rgb = C_BORDER_GRAY
    rubric_card.line.width = Pt(1.2)

    rtf = rubric_card.text_frame
    rtf.word_wrap = True
    rp0 = rtf.paragraphs[0]
    rp0.text = "STANDARDIZED EVALUATION CHECKLIST & MARK DISTRIBUTION"
    rp0.font.name = "Arial"
    rp0.font.size = Pt(14)
    rp0.font.bold = True
    rp0.font.color.rgb = C_NAVY_DARK

    rubric_items = [
        ("Conceptual Understanding & Formulation (30%)", "Correct identification of applicable theorems, formulas, historical events, or core principles."),
        ("Methodological Execution & Working Steps (40%)", "Logical sequence of derivations, intermediate algebraic steps, balanced chemical reactions, or source analysis."),
        ("Accuracy & Final Solution (20%)", "Precise numerical results, correct physical units, or coherent synthesis of concluding arguments."),
        ("Diagrams, Notation & Presentation (10%)", "Neatness, accurate diagram labeling, standard mathematical/scientific notation, and clarity.")
    ]

    for title, desc in rubric_items:
        p_t = rtf.add_paragraph()
        p_t.text = f"• {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(11.5)
        p_t.font.bold = True
        p_t.font.color.rgb = C_BLUE_ROYAL

        p_d = rtf.add_paragraph()
        p_d.text = f"   {desc}"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = C_NAVY_LIGHT

    _add_footer(slide_final, total_slides, total_slides, exam_title)

    # Save presentation
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    return output_path
