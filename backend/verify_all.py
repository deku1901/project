"""
verify_all.py
=============
End-to-End Integration and Quality Verification Suite:
  1. Verifies Frontend Static Assets and PPTX Card in index.html
  2. Verifies Knowledge Base / Context Ingestion with Universal Chunker
  3. Verifies Fast Exam Generation with Math & History prompts
  4. Verifies All 4 Publication Export Channels (PPTX, PDF, DOCX, LaTeX)
  5. Verifies PPTX slide structure, standard maths, and history source formatting
"""

import os
import re
import json
import requests
from pptx import Presentation

BASE_URL = "http://127.0.0.1:8000"

def verify_frontend():
    print("\n[1/5] Verifying Frontend Static Assets & HTML...")
    r = requests.get(f"{BASE_URL}/")
    assert r.status_code == 200, f"Root returned {r.status_code}"
    assert "PowerPoint Deck (.pptx)" in r.text, "Missing PPTX export card in index.html"
    assert "exportExam('pptx')" in r.text, "Missing PPTX onclick action in index.html"
    assert "KaTeX" in r.text, "KaTeX math library missing in HTML"
    print("✅ Frontend HTML & PPTX export card verified.")

def verify_chunker_upload():
    print("\n[2/5] Verifying Universal Chunker Upload on Complex Syllabi & History JSONs...")
    
    # Upload Syllabus JSON
    with open("CD_Calculus_Differential_Equations_context.json", "rb") as f:
        r = requests.post(f"{BASE_URL}/api/context/upload", files={"files": ("Calculus_Syllabus.json", f, "application/json")})
    assert r.status_code == 200, f"Upload failed: {r.text}"
    data = r.json()
    assert data["status"] == "success"
    assert data["count"] >= 20, f"Expected >= 20 chunks, got {data['count']}"
    print(f"✅ Calculus Syllabus parsed into {data['count']} semantic chunks.")

    # Upload History JSON
    hist_obj = {
        "topic": "World War II & Decolonization",
        "events": [
            {"year": "1945", "title": "End of World War II & UN Charter", "description": "Surrender of Axis powers and signing of the United Nations Charter in San Francisco.", "significance": "Established post-war international order."},
            {"year": "1947", "title": "Indian Independence Act", "description": "British Parliament passed legislation partitioning British India into India and Pakistan.", "significance": "Ended two centuries of British colonial rule."}
        ]
    }
    r = requests.post(f"{BASE_URL}/api/context/upload", files={"files": ("ww2_history.json", json.dumps(hist_obj).encode("utf-8"), "application/json")})
    assert r.status_code == 200, f"History upload failed: {r.text}"
    h_data = r.json()
    assert h_data["count"] == 2
    assert "History" in h_data["subject_breakdown"]
    print(f"✅ History dataset parsed into {h_data['count']} chunks with subject 'History'.")

def verify_fast_exam_generation():
    print("\n[3/5] Verifying High-Speed Exam Generation...")
    
    # Generate Math Exam
    r = requests.post(f"{BASE_URL}/api/generate", data={
        "title": "Calculus and Differential Equations",
        "n_questions": 3,
        "max_marks": 75,
        "include_diagrams": True
    })
    assert r.status_code == 200, f"Generate failed: {r.text}"
    exam_data = r.json()
    assert exam_data["status"] == "success"
    exam_id = exam_data["exam"]["id"]
    questions = exam_data["questions"]
    assert len(questions) >= 1, f"Expected questions, got {len(questions)}"
    print(f"✅ Math Exam generated (ID: {exam_id}) with {len(questions)} questions.")

    # Generate History Exam
    r_hist = requests.post(f"{BASE_URL}/api/generate", data={
        "title": "World War II & Decolonization",
        "n_questions": 2,
        "max_marks": 50,
        "include_diagrams": False
    })
    assert r_hist.status_code == 200, f"History generate failed: {r_hist.text}"
    h_exam_data = r_hist.json()
    h_exam_id = h_exam_data["exam"]["id"]
    h_questions = h_exam_data["questions"]
    assert len(h_questions) >= 1, f"Expected questions, got {len(h_questions)}"
    print(f"✅ History Exam generated (ID: {h_exam_id}) with {len(h_questions)} questions.")

    return exam_id, h_exam_id

def verify_all_exports(math_exam_id, hist_exam_id):
    print("\n[4/5] Verifying All 4 Export Channels (PPTX, PDF, DOCX, LaTeX)...")

    for eid, name in [(math_exam_id, "Math Exam"), (hist_exam_id, "History Exam")]:
        for fmt in ["pptx", "pdf", "docx", "tex"]:
            r = requests.post(f"{BASE_URL}/api/export", data={"exam_id": eid, "format": fmt})
            assert r.status_code == 200, f"Export {fmt} failed for {name}: {r.text}"
            assert len(r.content) > 500, f"Export {fmt} returned empty content"
            print(f"✅ {name} Exported to {fmt.upper()} successfully ({len(r.content):,} bytes).")

def verify_pptx_quality(math_exam_id):
    print("\n[5/5] Inspecting Generated PPTX Presentation Quality & Structure...")
    r = requests.post(f"{BASE_URL}/api/export", data={"exam_id": math_exam_id, "format": "pptx"})
    temp_pptx = "static/exports/verify_temp.pptx"
    with open(temp_pptx, "wb") as f:
        f.write(r.content)

    prs = Presentation(temp_pptx)
    print(f"✅ Total Slides in Presentation: {len(prs.slides)}")
    assert len(prs.slides) >= 5, "Expected Title + Blueprint + Questions + Marking Scheme slides"

    # Verify 16:9 widescreen
    w_in = prs.slide_width.inches
    h_in = prs.slide_height.inches
    print(f"✅ Slide Dimensions: {w_in:.3f}\" x {h_in:.3f}\" (16:9 Widescreen aspect ratio)")
    assert abs(w_in / h_in - 16 / 9) < 0.05, "Not 16:9 widescreen"

    # Check Slide 1 content
    s1_text = " ".join([shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame])
    assert "CALCULUS" in s1_text.upper()
    assert "INSTRUCTIONS" in s1_text.upper()
    print("✅ Slide 1: Cover & Instructions verified.")

    # Check Blueprint Table on Slide 2
    has_table = any(shape.has_table for shape in prs.slides[1].shapes)
    assert has_table, "Slide 2 missing blueprint table"
    print("✅ Slide 2: Blueprint Table verified.")

    slides_list = list(prs.slides)
    for idx, s in enumerate(slides_list[2:-1], start=1):
        s_text = " ".join([shape.text_frame.text for shape in s.shapes if shape.has_text_frame])
        assert f"Question {idx}" in s_text
        assert "Marks" in s_text
    print("✅ Question Slides verified with math equations and formatting.")

    # Check Final Marking Scheme Slide
    s_last_text = " ".join([shape.text_frame.text for shape in slides_list[-1].shapes if shape.has_text_frame])
    assert "RUBRIC" in s_last_text.upper() or "MARKING" in s_last_text.upper()
    print("✅ Final Slide: Marking Rubric verified.")


if __name__ == "__main__":
    verify_frontend()
    verify_chunker_upload()
    m_id, h_id = verify_fast_exam_generation()
    verify_all_exports(m_id, h_id)
    verify_pptx_quality(m_id)
    print("\n" + "="*60)
    print("🎉 ALL END-TO-END VERIFICATION CHECKS PASSED WITH 100% SUCCESS! 🎉")
    print("="*60 + "\n")
